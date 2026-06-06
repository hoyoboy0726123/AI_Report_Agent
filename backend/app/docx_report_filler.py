"""通用 docx 報告填圖引擎(參考 sn_ocr_tool 的「Agent 填圖」架構)。

核心原則(來自架構參考):
- 定位靠「文件結構」(表格/儲存格/blip rId),不靠 VLM 像素座標。
- 語意配對:範本槽位的「標籤(就近文字)」↔ 來源照片的「檔名語意」。
- 可靠副作用:**置換 docx 內部 media 位元組,document.xml 一字不動** → 圖片留在原位、不錯位、保留版面與尺寸框。
- 安全閘:數量/對應對不上就中止,不產半成品;不確定標 need_review。
- 驗證關卡:先把推導出的對應攤成清單給人看,確認後才寫檔。

流程:
  parse_template(docx)  → 槽位清單(每個圖 → 就近標籤 + rId + media 檔)
  parse_filenames(dir)  → 來源照片 → 語意標籤
  derive_mapping(...)   → 槽位 ↔ 照片(語意配對 + 順序消歧),含 need_review
  place_photos(...)     → 複製範本、置換 media 位元組、輸出 docx
"""

import os
import re
import shutil
import zipfile
from io import BytesIO

from app.image_mapper import _norm, _stem

NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# ============================================================
# A. 解析範本結構(可靠錨點)
# ============================================================

def _cell_blip_rids(cell):
    """回傳該儲存格內所有圖片的 r:embed rId(依出現順序)。"""
    xml = cell._tc.xml
    return re.findall(r'r:embed="([^"]+)"', xml)


def parse_template(docx_path: str) -> dict:
    """解析範本,回傳每個圖片槽位的結構錨點 + 就近標籤。

    回傳 {tables:[{table_index, slots:[{row,col,caption,rid,media}]}], total_slots, error?}
    槽位順序 = 文件閱讀順序(表 → 列 → 欄)。
    """
    from docx import Document
    if not docx_path or not os.path.isfile(docx_path):
        return {"error": f"檔案不存在: {docx_path}"}
    try:
        doc = Document(docx_path)
    except Exception as e:
        return {"error": f"無法開啟: {e}"}

    # rId → media partname(如 word/media/image1.jpeg)
    rid_to_media = {}
    try:
        for rid, rel in doc.part.rels.items():
            tgt = getattr(rel, "target_ref", "") or ""
            if "media/" in tgt:
                rid_to_media[rid] = "word/" + tgt if not tgt.startswith("word/") else tgt
    except Exception:
        pass

    tables_out = []
    total = 0
    for ti, t in enumerate(doc.tables):
        nrows = len(t.rows)
        ncols = len(t.columns)
        slots = []
        for ri, row in enumerate(t.rows):
            for ci, cell in enumerate(row.cells):
                rids = _cell_blip_rids(cell)
                if not rids:
                    continue
                # 就近標籤:優先「下一列同欄」的文字,其次同格文字、上一列同欄
                caption = ""
                if ri + 1 < nrows:
                    try:
                        caption = t.rows[ri + 1].cells[ci].text.strip()
                    except Exception:
                        caption = ""
                if not caption:
                    caption = cell.text.strip()
                if not caption and ri - 1 >= 0:
                    try:
                        caption = t.rows[ri - 1].cells[ci].text.strip()
                    except Exception:
                        caption = ""
                for rid in rids:
                    slots.append({
                        "row": ri, "col": ci,
                        "caption": caption.replace("\n", " ")[:60],
                        "rid": rid,
                        "media": rid_to_media.get(rid, ""),
                    })
                    total += 1
        if slots:
            tables_out.append({"table_index": ti, "slots": slots, "rows": nrows, "cols": ncols})
    return {"tables": tables_out, "total_slots": total,
            "photo_table_count": len(tables_out)}


def _pptx_iter_pictures(shapes):
    """遞迴 yield 投影片中的圖片形狀(含群組內)。"""
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            yield from _pptx_iter_pictures(sh.shapes)
        elif sh.shape_type == 13:  # PICTURE
            yield sh


def _pptx_text_shapes(shapes):
    out = []
    for sh in shapes:
        if sh.shape_type == 6:
            out.extend(_pptx_text_shapes(sh.shapes))
        elif sh.has_text_frame and (sh.text_frame.text or "").strip():
            out.append(sh)
    return out


def _center(sh):
    try:
        return (int(sh.left) + int(sh.width) / 2, int(sh.top) + int(sh.height) / 2)
    except Exception:
        return None


def _pic_media(slide, pic):
    """圖片形狀 → media partname(如 ppt/media/image1.png)。"""
    m = re.search(r'r:embed="([^"]+)"', pic._element.xml)
    if not m:
        return ""
    rid = m.group(1)
    try:
        part = slide.part.related_part(rid)
        return str(part.partname).lstrip("/")
    except Exception:
        return ""


def parse_template_pptx(pptx_path: str) -> dict:
    """PPTX 結構:每張投影片的圖片 + 最近文字當標籤。分組:每張投影片 = 一個樣本表。"""
    from pptx import Presentation
    if not pptx_path or not os.path.isfile(pptx_path):
        return {"error": f"檔案不存在: {pptx_path}"}
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"error": f"無法開啟: {e}"}

    tables_out, total = [], 0
    for si, slide in enumerate(prs.slides):
        pics = list(_pptx_iter_pictures(slide.shapes))
        if not pics:
            continue
        texts = _pptx_text_shapes(slide.shapes)
        slots = []
        for pic in pics:
            media = _pic_media(slide, pic)
            if not media:
                continue
            # 最近的文字當標籤(優先下方/上方近距離)
            caption = ""
            pc = _center(pic)
            if pc and texts:
                best, bestd = None, None
                for ts in texts:
                    tc = _center(ts)
                    if not tc:
                        continue
                    d = (tc[0] - pc[0]) ** 2 + (tc[1] - pc[1]) ** 2
                    if bestd is None or d < bestd:
                        bestd, best = d, ts
                if best is not None:
                    caption = (best.text_frame.text or "").strip().replace("\n", " ")[:60]
            slots.append({"row": si, "col": 0, "caption": caption, "rid": "", "media": media})
            total += 1
        if slots:
            tables_out.append({"table_index": si, "slots": slots, "rows": len(slots), "cols": 1})
    return {"tables": tables_out, "total_slots": total, "photo_table_count": len(tables_out)}


# ============================================================
# B. 解析來源檔名語意
# ============================================================

# 位置語意分類(類別, 在類別內的序號)。盡量用 token,不綁單一命名。
# 順序重要:corner / edge 必須在 top / bottom 之前判斷,
# 否則「Bottom corner」「Bottom edges」會被「bottom」先吃掉。
_CAT_PATTERNS = [
    ("before", r"before"),
    ("after", r"after"),
    ("appearance", r"appearance"),
    ("corner", r"corner"),
    ("edge", r"edge"),
    ("top", r"top"),
    ("bottom", r"bottom"),
]


def _classify_name(name: str):
    """檔名 → (category, index)。index 取檔名中該類別後面的數字(corner3→3),無則 0。"""
    low = name.lower()
    for cat, pat in _CAT_PATTERNS:
        m = re.search(pat, low)
        if m:
            # 抓類別關鍵字附近的數字
            num = 0
            tail = low[m.start():]
            dm = re.search(r"(\d+)", tail)
            if dm:
                num = int(dm.group(1))
            else:
                am = re.search(r"appearance(\d+)", low)
                if am:
                    num = int(am.group(1))
            return cat, num
    return "other", 0


def parse_filenames(folder: str) -> dict:
    """資料夾照片 → [{name, path, category, index}]。"""
    from app.image_mapper import list_images
    imgs = list_images(folder)
    out = []
    for im in imgs:
        cat, idx = _classify_name(im["name"])
        out.append({"name": im["name"], "path": im["path"], "category": cat, "index": idx})
    return {"files": out, "count": len(out)}


def _classify_caption(caption: str):
    low = caption.lower()
    for cat, pat in _CAT_PATTERNS:
        if re.search(pat, low):
            return cat
    return "other"


# ============================================================
# C. 推導對應 —— 通用作法:LLM 語意配對(不寫死關鍵字)
# ============================================================
# 結構解析給「可靠錨點 + 就近標籤」;這裡讓 LLM 理解「標籤語意 ↔ 檔名語意」自己配對,
# 因此換報告型態 / 換語言 / 換命名規則都不必改碼。
# 無 LLM 時退回 _heuristic_mapping(關鍵字,僅備援、不通用)。

_SEMANTIC_SYS = """你是報告填圖的「對應推導器」。會給你:
(1) 一個報告表格的「圖片槽位」清單,每個槽位有它在文件中的標籤文字(label)與出現順序;
(2) 一批「照片檔名」。
請純依語意,為每個槽位挑出最合適的照片檔名。規則:
- 標籤可能重複(例如多個 corner、多個 edge、多個視角),請用檔名中的編號或語意順序一一對應。
- 一個檔名最多只用一次;真的對不到就留空字串。
- 不要被特定命名綁住,依語意判斷(中英文、不同命名規則都要能處理)。
只回傳單一 JSON,不要任何其他文字:
{"mapping": [{"slot": 槽位編號(整數), "file": "對應檔名" 或 ""}]}"""


def derive_mapping_semantic(template_slots: list, files: list, llm, model: str) -> dict:
    """用 LLM 做語意配對(通用)。template_slots 需含 caption;files 為 [{name,path}]。"""
    import json as _json
    if llm is None or not model:
        return {"error": "未提供 LLM,無法語意配對", "need_review": True}
    slot_desc = [{"slot": i, "label": s.get("caption", ""), "order": i}
                 for i, s in enumerate(template_slots)]
    names = [f["name"] for f in files]
    user = _json.dumps({"圖片槽位": slot_desc, "照片檔名": names}, ensure_ascii=False)

    from app.agent.llm.base import Message
    try:
        resp = llm.chat([Message(role="system", text=_SEMANTIC_SYS),
                         Message(role="user", text=user)], model=model, tools=None)
        txt = (resp.text or "").strip()
        if txt.startswith("```"):
            txt = re.sub(r"^```\w*\n?", "", txt); txt = re.sub(r"\n?```\s*$", "", txt)
        data = _json.loads(txt)
        raw = data.get("mapping", data) if isinstance(data, dict) else data
    except Exception as e:
        return {"error": f"LLM 語意配對失敗: {e}", "need_review": True}

    name_to_file = {f["name"]: f for f in files}
    pairs, used = [], set()
    slot_filled = set()
    for entry in (raw or []):
        if not isinstance(entry, dict):
            continue
        si = entry.get("slot")
        fn = (entry.get("file") or "").strip()
        if si is None or fn == "" or fn not in name_to_file or fn in used:
            continue
        if not isinstance(si, int) or si < 0 or si >= len(template_slots):
            continue
        s = template_slots[si]
        f = name_to_file[fn]
        pairs.append({"caption": s.get("caption", ""), "rid": s["rid"], "media": s["media"],
                      "file": f["name"], "file_path": f["path"],
                      "row": s.get("row"), "col": s.get("col")})
        used.add(fn); slot_filled.add(si)
    unmatched_slots = [{"caption": s.get("caption", ""), "rid": s["rid"],
                        "row": s.get("row"), "col": s.get("col")}
                       for i, s in enumerate(template_slots) if i not in slot_filled]
    unmatched_files = [n for n in names if n not in used]
    return {"pairs": pairs, "unmatched_slots": unmatched_slots,
            "unmatched_files": unmatched_files,
            "need_review": bool(unmatched_slots or unmatched_files), "method": "llm-semantic"}


# ---- VLM 看圖配對(萬用:檔名無意義時靠照片內容 vs 範本示範圖) ----

def _extract_media(template_path: str, media_partnames: list) -> dict:
    """從 docx 抽出指定 media 檔到暫存,回 {partname: temp_path}。"""
    import tempfile
    out = {}
    want = set(media_partnames)
    tmpdir = tempfile.mkdtemp(prefix="tplmedia_")
    with zipfile.ZipFile(template_path, "r") as z:
        names = set(z.namelist())
        for pn in want:
            cand = pn if pn in names else pn.replace("word/", "", 1)
            if cand not in names:
                continue
            data = z.read(cand)
            ext = os.path.splitext(cand)[1] or ".jpg"
            p = os.path.join(tmpdir, os.path.basename(cand) if os.path.basename(cand) else f"m{len(out)}{ext}")
            with open(p, "wb") as f:
                f.write(data)
            out[pn] = p
    return out


def _thumb(path: str, max_dim: int = 512) -> str:
    """產生縮圖暫存(降 token / 加速視覺配對),回暫存路徑。失敗回原路徑。"""
    import tempfile
    from PIL import Image
    try:
        im = Image.open(path)
        if im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        w, h = im.size
        s = min(1.0, max_dim / float(max(w, h)))
        if s < 1.0:
            im = im.resize((int(w * s), int(h * s)), Image.LANCZOS)
        fd, p = tempfile.mkstemp(suffix=".jpg", prefix="thumb_")
        os.close(fd)
        im.save(p, format="JPEG", quality=80)
        return p
    except Exception:
        return path


_DESCRIBE_SYS = ("你是工程報告照片標註員。看一張照片,用最多 25 字描述它的「視角 + 主體 + 位置特徵」"
                 "(例:由上方俯視的筆電上蓋、機身右下角特寫、測試前整機外觀)。只回描述本身,不要前言。")


def _describe_photo(vlm, model, path, limiter=None) -> str:
    """VLM 看單張照片 → 一句內容描述(1 圖/次,不受多圖上限影響)。"""
    try:
        if limiter is not None:
            limiter.acquire()
        d = vlm.vision_complete(system=_DESCRIBE_SYS, user_text="描述這張照片。",
                                images=[_thumb(path)], model=model)
        return (d or "").strip().replace("\n", " ")[:60]
    except Exception:
        return ""


def _build_pairs_from_indexmap(template_slots, files, raw, descs=None):
    """從 [{slot,photo}] 索引對應建 pairs(共用)。"""
    pairs, used, filled = [], set(), set()
    for e in (raw or []):
        if not isinstance(e, dict):
            continue
        si, pj = e.get("slot"), e.get("photo")
        if not isinstance(si, int) or not isinstance(pj, int):
            continue
        if si < 0 or si >= len(template_slots) or pj < 0 or pj >= len(files) or pj in used:
            continue
        s, f = template_slots[si], files[pj]
        pr = {"caption": s.get("caption", ""), "rid": s["rid"], "media": s["media"],
              "file": f["name"], "file_path": f["path"], "row": s.get("row"), "col": s.get("col")}
        if descs:
            pr["photo_desc"] = descs[pj]
        pairs.append(pr); used.add(pj); filled.add(si)
    unmatched_slots = [{"caption": s.get("caption", ""), "rid": s["rid"]}
                       for i, s in enumerate(template_slots) if i not in filled]
    unmatched_files = [files[j]["name"] for j in range(len(files)) if j not in used]
    return {"pairs": pairs, "unmatched_slots": unmatched_slots,
            "unmatched_files": unmatched_files,
            "need_review": bool(unmatched_slots or unmatched_files)}


def _vlm_match_batched(template_slots, files, vlm, model):
    """一次批次:把所有照片(縮圖)+ 槽位標籤送一次 VLM,回 slot→photo 索引。

    少呼叫=少曝險(架構參考 2.11);需模型能吃多圖(gemini-2.5-flash 可,gemma 弱)。
    """
    import json as _json
    thumbs = [_thumb(f["path"]) for f in files]
    sys = ("你是報告填圖的視覺配對器。使用者會給你一組『待放照片』(依序就是接下來的圖片),"
           "以及一份『槽位清單』(每個槽位有標籤 label)。請看每張照片的實際內容(視角/主體/位置特徵),"
           "判斷它該放進哪個槽位。標籤重複(多個 corner/edge)時依視覺順序一一對應;"
           "一張照片最多用一次,對不到回 -1。只回 JSON:{\"mapping\":[{\"slot\":整數,\"photo\":整數或 -1}]}")
    slot_list = [{"slot": i, "label": s.get("caption", "")} for i, s in enumerate(template_slots)]
    photo_list = [{"photo": j, "name": files[j]["name"]} for j in range(len(files))]
    user = (f"槽位清單:\n{_json.dumps(slot_list, ensure_ascii=False)}\n\n"
            f"照片清單(共{len(files)}張,順序對應接下來的圖片):\n{_json.dumps(photo_list, ensure_ascii=False)}\n\n"
            "請依每張照片的實際內容配到最合適的槽位。")
    text = vlm.vision_complete(system=sys, user_text=user, images=thumbs, model=model)
    t = (text or "").strip()
    if t.startswith("```"):
        t = re.sub(r"^```\w*\n?", "", t); t = re.sub(r"\n?```\s*$", "", t)
    data = _json.loads(t)
    raw = data.get("mapping", data) if isinstance(data, dict) else data
    res = _build_pairs_from_indexmap(template_slots, files, raw)
    res["method"] = "vlm-batched"
    return res


def derive_mapping_vlm(template_slots: list, files: list, vlm, model: str,
                       template_path: str = "", strategy: str = "batched") -> dict:
    """視覺配對(萬用,適合檔名無語意)。

    strategy:
      - "batched":一次把所有照片+標籤送一次(少呼叫、快;需模型多圖能力強,如 gemini-2.5-flash)
      - "describe":逐張描述照片(1 圖/次,慢但弱模型如 gemma 單張較準)
      預設 batched;batched 失敗(如 gemma 多圖 400)會自動退回 describe。
    """
    if vlm is None or not model:
        return {"error": "未提供 VLM", "need_review": True}

    # 整批送(除非指定單張)。冷卻已集中在 client。
    if strategy != "describe":
        try:
            res = _vlm_match_batched(template_slots, files, vlm, model)
            if res.get("pairs"):
                return res
        except Exception:
            pass  # 多圖失敗 → 退回逐張描述

    import json as _json
    descs = [_describe_photo(vlm, model, f["path"]) for f in files]

    # 索引式語意配對(避免回傳檔名字串對不上的問題)
    sys = ("你是報告填圖配對器。會給你『槽位清單』(每個有標籤 label)與『照片清單』"
           "(每張有內容描述 content)。請依語意把每個槽位配到最合適的照片索引。"
           "標籤或描述重複時用順序/編號一一對應;一張照片最多用一次,對不到回 -1。"
           "只回 JSON:{\"mapping\":[{\"slot\":整數,\"photo\":整數或 -1}]}")
    slot_list = [{"slot": i, "label": s.get("caption", "")} for i, s in enumerate(template_slots)]
    photo_list = [{"photo": j, "content": descs[j]} for j in range(len(files))]
    user = _json.dumps({"槽位清單": slot_list, "照片清單": photo_list}, ensure_ascii=False)

    from app.agent.llm.base import Message
    try:
        resp = vlm.chat([Message(role="system", text=sys), Message(role="user", text=user)],
                        model=model, tools=None)
        txt = (resp.text or "").strip()
        if txt.startswith("```"):
            txt = re.sub(r"^```\w*\n?", "", txt); txt = re.sub(r"\n?```\s*$", "", txt)
        data = _json.loads(txt)
        raw = data.get("mapping", data) if isinstance(data, dict) else data
    except Exception as e:
        return {"error": f"VLM 配對失敗: {e}", "need_review": True, "descriptions": descs}

    pairs, used, filled = [], set(), set()
    for e in (raw or []):
        if not isinstance(e, dict):
            continue
        si, pj = e.get("slot"), e.get("photo")
        if not isinstance(si, int) or not isinstance(pj, int):
            continue
        if si < 0 or si >= len(template_slots) or pj < 0 or pj >= len(files) or pj in used:
            continue
        s, f = template_slots[si], files[pj]
        pairs.append({"caption": s.get("caption", ""), "rid": s["rid"], "media": s["media"],
                      "file": f["name"], "file_path": f["path"], "row": s.get("row"),
                      "col": s.get("col"), "photo_desc": descs[pj]})
        used.add(pj); filled.add(si)
    unmatched_slots = [{"caption": s.get("caption", ""), "rid": s["rid"]}
                       for i, s in enumerate(template_slots) if i not in filled]
    unmatched_files = [files[j]["name"] for j in range(len(files)) if j not in used]
    return {"pairs": pairs, "unmatched_slots": unmatched_slots,
            "unmatched_files": unmatched_files,
            "need_review": bool(unmatched_slots or unmatched_files),
            "method": "vlm-describe"}


def derive_mapping_semantic_batched(samples: list, llm, model: str) -> dict:
    """合併呼叫:一次處理多個樣本的「標籤↔檔名」語意配對(少呼叫=少曝險,架構參考 2.11)。

    samples = [{"key": 樣本鍵, "slots": [...], "files": [{name,path}]}]
    回 {key: {pairs, unmatched_slots, unmatched_files, need_review, method}}。
    用「樣本索引 + 槽位索引 + 照片索引」回傳,避免檔名字串對不上。
    """
    import json as _json
    if llm is None or not model:
        return {"error": "未提供 LLM"}
    payload = []
    for si, s in enumerate(samples):
        payload.append({
            "sample": si, "key": s["key"],
            "slots": [{"slot": i, "label": sl.get("caption", "")} for i, sl in enumerate(s["slots"])],
            "photos": [{"photo": j, "name": f["name"]} for j, f in enumerate(s["files"])],
        })
    sys = ("你是報告填圖配對器。會給你多個『樣本』,每個樣本有自己的槽位(label)與照片(name)。"
           "請為每個樣本、每個槽位,配一張『同樣本內』最合適的照片(依語意:標籤↔檔名)。"
           "標籤重複(多個 corner/edge)用檔名編號或順序一一對應;同樣本內一張照片最多用一次,對不到回 -1。"
           "只回 JSON:{\"samples\":[{\"sample\":整數,\"mapping\":[{\"slot\":整數,\"photo\":整數或 -1}]}]}")
    user = _json.dumps(payload, ensure_ascii=False)

    from app.agent.llm.base import Message
    try:
        resp = llm.chat([Message(role="system", text=sys), Message(role="user", text=user)],
                        model=model, tools=None)
        txt = (resp.text or "").strip()
        if txt.startswith("```"):
            txt = re.sub(r"^```\w*\n?", "", txt); txt = re.sub(r"\n?```\s*$", "", txt)
        data = _json.loads(txt)
        sample_maps = data.get("samples", []) if isinstance(data, dict) else data
    except Exception as e:
        return {"error": f"合併語意配對失敗: {e}"}

    by_idx = {}
    for entry in (sample_maps or []):
        if isinstance(entry, dict) and isinstance(entry.get("sample"), int):
            by_idx[entry["sample"]] = entry.get("mapping", [])

    out = {}
    for si, s in enumerate(samples):
        raw = by_idx.get(si, [])
        res = _build_pairs_from_indexmap(s["slots"], s["files"], raw)
        res["method"] = "text-batched"
        out[s["key"]] = res
    return out


def _slot_pair(s, f):
    return {"caption": s.get("caption", ""), "rid": s.get("rid", ""), "media": s.get("media", ""),
            "file": f["name"], "file_path": f["path"], "row": s.get("row"), "col": s.get("col")}


def _exact_prematch(slots, files):
    """確定性預配:標籤正規化後 == 檔名 stem(或唯一包含)→ 直接配,不勞煩模型。
    回 (pairs, remaining_slots, remaining_files)。"""
    norm_files = [(j, _norm(_stem(f["name"]))) for j, f in enumerate(files)]
    used, matched = set(), set()
    pairs = []
    # pass1 精確相等且唯一
    for i, s in enumerate(slots):
        cap = _norm(s.get("caption", ""))
        if not cap:
            continue
        cands = [j for j, nf in norm_files if j not in used and nf == cap]
        if len(cands) == 1:
            pairs.append(_slot_pair(s, files[cands[0]])); used.add(cands[0]); matched.add(i)
    # pass2 唯一包含(caption ⊆ stem 或 stem ⊆ caption)
    for i, s in enumerate(slots):
        if i in matched:
            continue
        cap = _norm(s.get("caption", ""))
        if not cap:
            continue
        cands = [j for j, nf in norm_files if j not in used and nf and (cap in nf or nf in cap)]
        if len(cands) == 1:
            pairs.append(_slot_pair(s, files[cands[0]])); used.add(cands[0]); matched.add(i)
    rem_slots = [s for i, s in enumerate(slots) if i not in matched]
    rem_files = [f for j, f in enumerate(files) if j not in used]
    return pairs, rem_slots, rem_files


def derive_mapping(template_slots: list, files: list, llm=None, model: str = "",
                   vlm=None, vlm_model: str = "", template_path: str = "",
                   mode: str = "auto") -> dict:
    """對外入口:先做確定性精確預配,剩下對不齊的再依 mode 用模型配。

    mode: text(LLM 檔名↔標籤) / vlm(視覺) / auto(text→VLM 補強,預設)。
    """
    # 確定性預配(完全相同的標籤↔檔名不該交給模型猜)
    pre_pairs, rem_slots, rem_files = _exact_prematch(template_slots, files)
    if not rem_slots or not rem_files:
        return {"pairs": pre_pairs,
                "unmatched_slots": [{"caption": s.get("caption", ""), "rid": s.get("rid", "")} for s in rem_slots],
                "unmatched_files": [f["name"] for f in rem_files],
                "need_review": bool(rem_slots or rem_files), "method": "exact"}
    core = _match_core(rem_slots, rem_files, llm, model, vlm, vlm_model, template_path, mode)
    core["pairs"] = pre_pairs + core.get("pairs", [])
    if pre_pairs:
        core["method"] = "exact+" + core.get("method", "")
    return core


def _match_core(template_slots: list, files: list, llm=None, model: str = "",
                vlm=None, vlm_model: str = "", template_path: str = "",
                mode: str = "auto") -> dict:
    if mode == "vlm":
        if vlm is not None and vlm_model and template_path:
            r = derive_mapping_vlm(template_slots, files, vlm, vlm_model, template_path)
            if "error" not in r:
                return r
        return _heuristic_mapping(template_slots, files)

    # text first
    res = None
    if llm is not None and model:
        res = derive_mapping_semantic(template_slots, files, llm, model)
    if res is None or "error" in res:
        res = _heuristic_mapping(template_slots, files)

    # auto:對不齊就用 VLM 針對「未配槽位 + 未配檔案」補強
    if mode == "auto" and res.get("need_review") and vlm is not None and vlm_model and template_path:
        un_slots = [s for s in template_slots
                    if any(u["rid"] == s["rid"] for u in res.get("unmatched_slots", []))]
        un_names = set(res.get("unmatched_files", []))
        un_files = [f for f in files if f["name"] in un_names]
        if un_slots and un_files:
            vres = derive_mapping_vlm(un_slots, un_files, vlm, vlm_model, template_path)
            if "error" not in vres and vres.get("pairs"):
                res["pairs"].extend(vres["pairs"])
                done_rids = {p["rid"] for p in res["pairs"]}
                used_files = {p["file"] for p in res["pairs"]}
                res["unmatched_slots"] = [u for u in res["unmatched_slots"] if u["rid"] not in done_rids]
                res["unmatched_files"] = [n for n in res["unmatched_files"] if n not in used_files]
                res["need_review"] = bool(res["unmatched_slots"] or res["unmatched_files"])
                res["method"] = "text+vlm"
    return res


def _heuristic_mapping(template_slots: list, files: list) -> dict:
    """單一樣本:把一個照片表的 slots ↔ 一個資料夾的 files 配對。

    作法:各自分類;同類別內依「槽位閱讀順序」對「檔名序號排序」一一配對。
    回傳 {pairs:[{slot, file, caption, category}], unmatched_slots, unmatched_files, need_review}
    """
    # 槽位分類(保留閱讀順序)
    slots_by_cat = {}
    for s in template_slots:
        cat = _classify_caption(s["caption"])
        slots_by_cat.setdefault(cat, []).append(s)
    # 檔案分類(同類別內依 index 排序,確保 corner1<2<3<4)
    files_by_cat = {}
    for f in files:
        files_by_cat.setdefault(f["category"], []).append(f)
    for cat in files_by_cat:
        files_by_cat[cat].sort(key=lambda x: x["index"])

    pairs = []
    unmatched_slots = []
    used_files = set()
    for cat, slots in slots_by_cat.items():
        fs = files_by_cat.get(cat, [])
        for i, s in enumerate(slots):
            if i < len(fs):
                f = fs[i]
                pairs.append({"caption": s["caption"], "category": cat,
                              "rid": s["rid"], "media": s["media"],
                              "file": f["name"], "file_path": f["path"],
                              "row": s["row"], "col": s["col"]})
                used_files.add(f["path"])
            else:
                unmatched_slots.append({"caption": s["caption"], "category": cat,
                                        "rid": s["rid"], "row": s["row"], "col": s["col"]})
    unmatched_files = [f["name"] for f in files if f["path"] not in used_files]
    need_review = bool(unmatched_slots or unmatched_files)
    return {"pairs": pairs, "unmatched_slots": unmatched_slots,
            "unmatched_files": unmatched_files, "need_review": need_review}


# ============================================================
# D. 可靠副作用:置換 media 位元組(XML 不動)
# ============================================================

def _encode_photo(path: str, max_dim: int = 1600, quality: int = 85) -> bytes:
    """讀照片 → 縮到 max_dim 內 → 輸出 jpeg bytes(控檔案大小,避免百 MB)。"""
    from PIL import Image
    im = Image.open(path)
    if im.mode not in ("RGB", "L"):
        im = im.convert("RGB")
    w, h = im.size
    scale = min(1.0, max_dim / float(max(w, h)))
    if scale < 1.0:
        im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
    buf = BytesIO()
    im.save(buf, format="JPEG", quality=quality)
    return buf.getvalue()


def place_photos(template_path: str, media_to_photo: dict, output_path: str,
                 max_dim: int = 1600, quality: int = 85) -> dict:
    """複製範本、把指定 media 檔的位元組換成新照片,輸出 docx。document.xml 完全不動。

    media_to_photo = {word/media/imageN.jpeg: 新照片路徑}
    """
    if not os.path.isfile(template_path):
        return {"error": f"範本不存在: {template_path}"}
    if not media_to_photo:
        return {"error": "沒有要置換的 media(對應為空)"}

    # 預先編碼所有新照片
    encoded = {}
    for media, photo in media_to_photo.items():
        if not photo or not os.path.isfile(photo):
            return {"error": f"來源照片不存在: {photo}"}
        try:
            encoded[media] = _encode_photo(photo, max_dim, quality)
        except Exception as e:
            return {"error": f"編碼照片失敗 {os.path.basename(photo)}: {e}"}

    # 依「media 檔名(basename)」比對,docx(word/media)/pptx(ppt/media)/xlsx(xl/media)通吃。
    # 單一檔內 media basename 唯一,安全。
    encoded_by_base = {os.path.basename(k): v for k, v in encoded.items()}

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    try:
        zin = zipfile.ZipFile(template_path, "r")
        zout = zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED)
        replaced = 0
        for item in zin.infolist():
            data = zin.read(item.filename)
            key = item.filename
            base = os.path.basename(key)
            if key in encoded:                         # 完整 partname 命中
                data = encoded[key]; replaced += 1
            elif "/media/" in key and base in encoded_by_base:   # 依 media 檔名命中(通用)
                data = encoded_by_base[base]; replaced += 1
            zout.writestr(item, data)
        zin.close(); zout.close()
    except Exception as e:
        return {"error": f"寫檔失敗: {e}"}
    return {"ok": True, "output": output_path, "replaced": replaced,
            "expected": len(media_to_photo)}

"""PPTX 範本標籤模式:讀含 {{tag}} 的 .pptx 範本 → 依來源 Excel 每列產出新 .pptx。

對照:
- ReportGenerator       → Word 範本 + Excel → 多份 .docx
- ExcelReportGenerator  → Excel 範本 + Excel → 多份 .xlsx
- PptxReportGenerator   → PPTX 範本 + Excel → 多份 .pptx

填充規則:
- 掃所有投影片的「文字框 / 表格儲存格 / 群組內形狀」中的 {{tag}},以該列資料替換。
- run 級替換(保留格式);{{tag}} 跨 run 時退回段落級重寫。
- 若某形狀的文字「整個就是單一 {{tag}}」且值為圖片路徑 → 以該形狀的位置/大小插入圖片並清空文字。
"""

import os
import re

import pandas as pd

from app.config import DEFAULT_OUTPUT_DIR, DEFAULT_HEADER_ROW, IMAGE_EXTENSIONS
from app.filename import render_filename

DEFAULT_PPTX_FILENAME_TEMPLATE = "報告_{index}.pptx"
TAG_RE = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")


def _iter_text_frames_and_shapes(shapes):
    """遞迴走訪所有形狀,yield (shape, text_frame_or_None);表格另外處理。"""
    for sh in shapes:
        # 群組 → 遞迴
        if sh.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_text_frames_and_shapes(sh.shapes)
            continue
        yield sh


def _scan_tags_in_pptx(prs) -> set:
    found = set()
    for slide in prs.slides:
        for sh in _iter_text_frames_and_shapes(slide.shapes):
            if sh.has_text_frame:
                for m in TAG_RE.findall(sh.text_frame.text or ""):
                    found.add(m.strip())
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        for m in TAG_RE.findall(cell.text or ""):
                            found.add(m.strip())
    return found


def _replace_in_text_frame(tf, resolve):
    """對一個 text_frame 內所有段落做 {{tag}} 替換(run 級優先,跨 run 退段落級)。"""
    for para in tf.paragraphs:
        runs = para.runs
        full = "".join(r.text for r in runs)
        if not TAG_RE.search(full):
            continue
        # 先試 run 級(保留格式)
        changed_run = False
        for r in runs:
            if TAG_RE.search(r.text or ""):
                r.text = TAG_RE.sub(lambda m: resolve(m.group(1).strip(), m.group(0)), r.text)
                changed_run = True
        # 若 tag 跨 run(整段仍有殘留)→ 段落級重寫
        new_full = "".join(r.text for r in para.runs)
        if TAG_RE.search(new_full):
            replaced = TAG_RE.sub(lambda m: resolve(m.group(1).strip(), m.group(0)), new_full)
            for r in para.runs:
                r.text = ""
            if para.runs:
                para.runs[0].text = replaced
        _ = changed_run


class PptxReportGenerator:
    def __init__(self, template_path: str, excel_path: str,
                 output_dir: str = DEFAULT_OUTPUT_DIR, sheet_name: str = "",
                 header_row: int = DEFAULT_HEADER_ROW,
                 filename_template: str = DEFAULT_PPTX_FILENAME_TEMPLATE):
        self.template_path = template_path
        self.excel_path = excel_path
        self.output_dir = output_dir
        self.sheet_name = sheet_name if sheet_name else 0
        self.header_row = max(1, int(header_row or 1))
        self.filename_template = filename_template or DEFAULT_PPTX_FILENAME_TEMPLATE

    # ---------- 掃描 / 驗證 ----------

    def template_variables(self) -> set:
        from pptx import Presentation
        if not os.path.isfile(self.template_path):
            raise FileNotFoundError(f"PPTX 範本不存在: {self.template_path}")
        return _scan_tags_in_pptx(Presentation(self.template_path))

    def list_sheets(self):
        return pd.ExcelFile(self.excel_path).sheet_names

    def _read_dataframe(self):
        return pd.read_excel(self.excel_path, sheet_name=self.sheet_name,
                             header=self.header_row - 1)

    def validate(self):
        template_vars = self.template_variables()
        df = self._read_dataframe()
        excel_cols = {str(c) for c in df.columns}
        missing = template_vars - excel_cols
        extra = excel_cols - template_vars
        return missing, extra

    # ---------- 值解析 ----------

    @staticmethod
    def _is_image_path(s):
        return (isinstance(s, str) and s.lower().endswith(IMAGE_EXTENSIONS)
                and os.path.isfile(s))

    def _resolve_factory(self, row_dict):
        def resolve(key, original):
            if key not in row_dict:
                return original  # 找不到 → 保留原 tag
            v = row_dict[key]
            try:
                if pd.isna(v):
                    return ""
            except (TypeError, ValueError):
                pass
            return str(v)
        return resolve

    def _maybe_image_shape(self, slide, shape, row_dict):
        """形狀文字整個=單一 {{tag}} 且值為圖片 → 就地插圖、清文字。回 True 表示已處理。"""
        if not shape.has_text_frame:
            return False
        text = (shape.text_frame.text or "").strip()
        m = TAG_RE.fullmatch(text)
        if not m:
            return False
        key = m.group(1).strip()
        v = row_dict.get(key)
        if not self._is_image_path(v):
            return False
        try:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes.add_picture(v, left, top, width=width, height=height)
            shape.text_frame.text = ""  # 清掉佔位文字
            return True
        except Exception:
            return False

    # ---------- 產出 ----------

    def generate_iter(self, cancel_event=None):
        from pptx import Presentation
        df = self._read_dataframe()
        os.makedirs(self.output_dir, exist_ok=True)
        total = len(df)
        produced = 0
        for index, row in df.iterrows():
            if cancel_event is not None and cancel_event.is_set():
                return
            prs = Presentation(self.template_path)
            row_dict = row.to_dict()
            resolve = self._resolve_factory(row_dict)
            for slide in prs.slides:
                for sh in list(_iter_text_frames_and_shapes(slide.shapes)):
                    # 圖片優先(整格單一 tag 且值為圖)
                    if self._maybe_image_shape(slide, sh, row_dict):
                        continue
                    if sh.has_text_frame:
                        _replace_in_text_frame(sh.text_frame, resolve)
                    if sh.has_table:
                        for r in sh.table.rows:
                            for cell in r.cells:
                                _replace_in_text_frame(cell.text_frame, resolve)
            filename = render_filename(self.filename_template, row_dict, index + 1,
                                       default_ext=".pptx")
            saved = os.path.join(self.output_dir, filename)
            prs.save(saved)
            produced += 1
            yield produced, total, saved, row_dict

    def generate(self, progress_callback=None, cancel_event=None):
        produced = total = 0
        for prod, tot, _, _ in self.generate_iter(cancel_event=cancel_event):
            produced, total = prod, tot
            if progress_callback:
                progress_callback(prod, tot)
        return produced, total


# ============================================================
# 視覺自動配圖(PPT 無標籤):把資料夾照片依內容/情境貼到對的投影片
# ============================================================

def _slide_caption(slide):
    """取投影片的代表文字(標題優先,否則最長的一段文字)當情境標籤。"""
    title = ""
    try:
        if slide.shapes.title and (slide.shapes.title.text or "").strip():
            title = slide.shapes.title.text.strip()
    except Exception:
        title = ""
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame and (sh.text_frame.text or "").strip():
            texts.append(sh.text_frame.text.strip().replace("\n", " "))
    body = max(texts, key=len) if texts else ""
    cap = (title + " " + body).strip() if title else body
    return cap[:80]


def _place_on_slide(slide, img_path, width_in=4.5):
    """把圖片放進投影片:優先空的圖片版面配置區,否則加在預設位置。"""
    from pptx.util import Inches
    for ph in list(slide.placeholders):
        try:
            if ph.placeholder_format.type == 18:  # PP_PLACEHOLDER.PICTURE
                ph.insert_picture(img_path)
                return "placeholder"
        except Exception:
            continue
    slide.shapes.add_picture(img_path, Inches(1.5), Inches(2.6), width=Inches(width_in))
    return "added"


def auto_place_images_pptx(vlm, model, pptx_path, image_folder, output_path,
                           llm=None, llm_model="", mode="auto", rpm=0, width_in=4.5):
    """PPT 視覺自動配圖。回傳結果摘要(含 placements)。"""
    import os as _os
    from pptx import Presentation
    from app.docx_report_filler import parse_filenames, derive_mapping

    if not pptx_path or not _os.path.isfile(pptx_path):
        return {"error": f"PPTX 不存在: {pptx_path}"}
    files = parse_filenames(image_folder)["files"]
    if not files:
        return {"error": f"資料夾無圖片: {image_folder}"}

    prs = Presentation(pptx_path)
    # 每張投影片 = 一個槽位(caption = 投影片文字)
    slots = [{"caption": _slide_caption(s), "row": i, "rid": "", "media": ""}
             for i, s in enumerate(prs.slides)]
    if not slots:
        return {"error": "簡報沒有投影片"}

    m = derive_mapping(slots, files, llm=llm, model=llm_model,
                       vlm=vlm, vlm_model=model, template_path=pptx_path, mode=mode)
    pairs = m.get("pairs", [])
    if not pairs:
        return {"error": "未能配對任何照片到投影片", "need_review": True,
                "unmatched_files": m.get("unmatched_files", [])}

    # 插入
    slides = list(prs.slides)
    placed = []
    for p in pairs:
        si = p.get("row")
        if si is None or si < 0 or si >= len(slides):
            continue
        where = _place_on_slide(slides[si], p["file_path"], width_in=width_in)
        placed.append({"slide": si + 1, "file": p["file"], "caption": p["caption"], "where": where})

    if not output_path:
        base, ext = _os.path.splitext(pptx_path)
        output_path = f"{base}_filled{ext}"
    _os.makedirs(_os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return {"ok": True, "output": output_path, "placed": placed,
            "placed_count": len(placed), "total_images": len(files),
            "unmatched_files": m.get("unmatched_files", []), "method": m.get("method", "")}


# ============================================================
# 範本共用靜態圖(PPT):把文字=單一 {{欄位}} 的形狀就地換成圖片
# ============================================================

def replace_tag_with_image_pptx(pptx_path, field, image_path, output_path=None):
    """找文字整個是 {{field}} 的形狀,用其位置/大小插圖、清文字。回 {replaced}。"""
    import os as _os
    from pptx import Presentation
    if not _os.path.isfile(pptx_path):
        return {"error": f"PPTX 不存在: {pptx_path}"}
    if not image_path or not _os.path.isfile(image_path):
        return {"error": f"圖片不存在: {image_path}"}
    pat = re.compile(r"\{\{\s*" + re.escape(field) + r"\s*\}\}")
    prs = Presentation(pptx_path)
    replaced = 0
    for slide in prs.slides:
        for sh in list(_iter_text_frames_and_shapes(slide.shapes)):
            if sh.has_text_frame and pat.fullmatch((sh.text_frame.text or "").strip()):
                try:
                    l, t, w, h = sh.left, sh.top, sh.width, sh.height
                    slide.shapes.add_picture(image_path, l, t, width=w, height=h)
                    sh.text_frame.text = ""
                    replaced += 1
                except Exception as e:
                    return {"error": f"插入圖片失敗: {e}"}
    if replaced:
        prs.save(output_path or pptx_path)
    return {"replaced": replaced, "field": field}


# ============================================================
# 無標註文字填寫(PPT):把 Excel 一列資料依語意填到投影片標籤後
# ============================================================

def auto_fill_text_pptx(llm, model, pptx_path, data, output_path):
    """PPT 無 {{標籤}} 時,讓 LLM 判斷每個欄位值該接在哪張投影片的哪段標籤文字後,寫入。"""
    import os as _os
    import json as _json
    from pptx import Presentation
    from app.agent.llm.base import Message

    if not pptx_path or not _os.path.isfile(pptx_path):
        return {"error": f"PPTX 不存在: {pptx_path}"}
    data = {str(k): ("" if v is None else str(v)) for k, v in (data or {}).items()}
    data = {k: v for k, v in data.items() if v.strip()}
    if not data:
        return {"error": "沒有可填的資料"}

    prs = Presentation(pptx_path)
    # 蒐集每張投影片的文字形狀(供定位)
    slide_texts = []
    shape_index = {}  # (slide_idx, shape_id) -> shape
    for si, slide in enumerate(prs.slides):
        items = []
        for sh in slide.shapes:
            if sh.has_text_frame and (sh.text_frame.text or "").strip():
                items.append({"sid": sh.shape_id, "text": sh.text_frame.text.strip().replace("\n", " ")[:80]})
                shape_index[(si, sh.shape_id)] = sh
        if items:
            slide_texts.append({"slide": si, "texts": items})
    if not slide_texts:
        return {"error": "簡報沒有可定位的文字"}

    sys = ("你是簡報填寫助理。會給你一份『欄位:值』資料,與各投影片現有的文字形狀(每個有 slide、sid、text)。"
           "請判斷每個欄位的值該接在哪個文字形狀後面(通常是標籤,如「客戶:」)。"
           "只回 JSON:{\"fills\":[{\"field\":\"欄位\",\"slide\":整數,\"sid\":整數}]};對不到的欄位略過。")
    user = _json.dumps({"資料": data, "投影片文字": slide_texts}, ensure_ascii=False)
    try:
        resp = llm.chat([Message(role="system", text=sys), Message(role="user", text=user)],
                        model=model, tools=None)
        txt = (resp.text or "").strip()
        if txt.startswith("```"):
            import re as _re
            txt = _re.sub(r"^```\w*\n?", "", txt); txt = _re.sub(r"\n?```\s*$", "", txt)
        fills = _json.loads(txt).get("fills", [])
    except Exception as e:
        return {"error": f"LLM 配對失敗: {e}"}

    filled, failed = [], []
    for f in fills:
        field = f.get("field"); si = f.get("slide"); sid = f.get("sid")
        if field not in data:
            continue
        sh = shape_index.get((si, sid))
        if sh is None:
            failed.append({"field": field, "reason": "找不到目標文字形狀"})
            continue
        val = data[field]
        tf = sh.text_frame
        cur = tf.text
        # 標籤結尾是冒號就直接接,否則空格分隔
        sep = "" if cur.rstrip().endswith((":", "：")) else " "
        try:
            para = tf.paragraphs[-1]
            r = para.add_run()
            r.text = sep + val
            filled.append({"field": field, "slide": (si or 0) + 1, "anchor": cur[:20]})
        except Exception as e:
            failed.append({"field": field, "reason": str(e)})

    if not output_path:
        base, ext = _os.path.splitext(pptx_path)
        output_path = f"{base}_filled{ext}"
    _os.makedirs(_os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return {"ok": True, "output": output_path, "filled": filled, "fill_failed": failed,
            "filled_count": len(filled), "total_fields": len(data)}
